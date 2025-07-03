import io
import logging
import os
import subprocess
import tempfile
import pandas as pd
from pptx import Presentation
from unstructured.partition.auto import partition

logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_content: bytes) -> str:
    logger.info("🧾 Starting PDF text extraction using Unstructured...")
    try:
        file_like = io.BytesIO(file_content)
        elements = partition(
            file=file_like,
            strategy="hi_res",
            pdf_image_dpi=300,
            languages=["eng"],
        )
        text = " ".join(str(element) for element in elements if str(element).strip())
        logger.info(f" Extracted PDF text length (hi_res): {len(text)}")
        if text.strip():
            return text
    except Exception as e:
        logger.warning(f" Unstructured PDF extraction failed: {e}")

    logger.info(" Falling back to OCR for PDF using tesseract...")
    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "temp.pdf")
        with open(pdf_path, "wb") as f:
            f.write(file_content)

        images_dir = os.path.join(temp_dir, "images")
        os.makedirs(images_dir, exist_ok=True)

        try:
            subprocess.run(["pdftoppm", "-jpeg", pdf_path, os.path.join(images_dir, "page")], check=True)
            extracted_text = ""
            for img_file in sorted(os.listdir(images_dir)):
                img_path = os.path.join(images_dir, img_file)
                ocr_text = subprocess.check_output(["tesseract", img_path, "stdout"], encoding="utf-8")
                extracted_text += ocr_text + "\n"
            logger.info(f" Extracted PDF text length (OCR): {len(extracted_text)}")
            return extracted_text.strip()
        except Exception as e:
            logger.error(f" OCR extraction failed: {e}")
            return ""

def parse_document(file_content: bytes, content_type: str) -> str:
    logger.info(f" Parsing document of type: {content_type}")
    try:
        file_like = io.BytesIO(file_content)

        if content_type == "application/pdf":
            return extract_text_from_pdf(file_content)

        elif content_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            logger.info(" Extracting text from PPTX...")
            prs = Presentation(file_like)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            output = " ".join(text).strip()
            logger.info(f" Extracted PPTX text length: {len(output)}")
            return output

        elif content_type == "text/csv":
            logger.info(" Extracting text from CSV...")
            content_str = file_content.decode("utf-8", errors="ignore")
            df = pd.read_csv(io.StringIO(content_str))
            output = df.to_string()
            logger.info(f" Parsed CSV text length: {len(output)}")
            return output

        elif content_type == "text/plain":
            logger.info(" Extracting text from plain text file...")
            output = file_content.decode("utf-8", errors="ignore").strip()
            logger.info(f" Extracted TXT text length: {len(output)}")
            return output

        else:
            raise ValueError(f" Unsupported file type: {content_type}")

    except Exception as e:
        logger.error(f" Error parsing document: {e}", exc_info=True)
        return ""
