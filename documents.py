from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File
from fastapi.security import OAuth2PasswordBearer
from sqlalchemy.orm import Session
from ..database import get_db
from ..models import Document, User
from ..schemas import DocumentCreate, DocumentResponse
from ..config import S3_ENDPOINT_URL, S3_ACCESS_KEY, S3_SECRET_KEY, S3_BUCKET, SECRET_KEY, ALGORITHM
from jose import jwt, JWTError
import boto3
import uuid
import io
from pptx import Presentation
import pandas as pd
import logging
from .rag import index_document
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="auth/login")

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")
    except JWTError:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")
    
    user = db.query(User).filter(User.username == username).first()
    if user is None:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="User not found")
    return user

router = APIRouter(tags=["documents"])
__all__ = ["router"]

s3_client = boto3.client(
    "s3",
    endpoint_url=S3_ENDPOINT_URL,
    aws_access_key_id=S3_ACCESS_KEY,
    aws_secret_access_key=S3_SECRET_KEY
)

def parse_document(file_content, content_type):
    try:
        file_like = io.BytesIO(file_content)
        if content_type == "application/pdf":
            from pikepdf import Pdf
            text = []
            with Pdf.open(file_like) as pdf:
                for page in pdf.pages:
                    extracted_text = page.get("/Contents", default="").get_text() or ""
                    if extracted_text:
                        text.append(extracted_text)
                    else:
                        # Fallback to extracting text from page objects if direct content fails
                        if "/XObject" in page:
                            for xobj in page["/XObject"].values():
                                if "/Subtype" in xobj and xobj["/Subtype"] == "/Form":
                                    form_text = xobj.get_text() or ""
                                    if form_text:
                                        text.append(form_text)
            return " ".join(text) if text else ""
        elif content_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            prs = Presentation(file_like)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text)
        elif content_type == "text/csv":
            content_str = file_content.decode('utf-8')
            logger.info(f"CSV content: {content_str}")
            df = pd.read_csv(io.StringIO(content_str))
            logger.info(f"Parsed DataFrame: {df.to_string()}")
            return df.to_string()
        else:
            raise ValueError(f"Unsupported file type: {content_type}")
    except Exception as e:
        logger.error(f"Error parsing document: {str(e)}")
        return ""

@router.post("/upload", response_model=DocumentResponse)
async def upload_document(file: UploadFile = File(...), db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    logger.info(f"Received file with content_type: {file.content_type}")
    allowed_types = [
        "application/pdf",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "text/csv",
        "application/octet-stream"
    ]
    if file.content_type not in allowed_types:
        logger.error(f"Unsupported file type: {file.content_type}")
        raise HTTPException(status_code=400, detail="Unsupported file type")

    file_extension = file.filename.split(".")[-1].lower()
    unique_filename = f"{uuid.uuid4()}.{file_extension}"

    file_content = await file.read()
    file.file.seek(0)

    try:
        s3_client.upload_fileobj(io.BytesIO(file_content), S3_BUCKET, unique_filename)
    except Exception as e:
        logger.error(f"Failed to upload file to S3: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to upload file: {str(e)}")

    if file.content_type == "application/octet-stream":
        if file_extension == "pptx":
            content_type_for_parsing = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif file_extension == "ppt":
            content_type_for_parsing = "application/vnd.ms-powerpoint"
        elif file_extension == "csv":
            content_type_for_parsing = "text/csv"
        elif file_extension == "pdf":
            content_type_for_parsing = "application/pdf"
        else:
            logger.error(f"Unsupported file extension: {file_extension}")
            raise HTTPException(status_code=400, detail="Unsupported file extension")
    else:
        content_type_for_parsing = file.content_type

    content = parse_document(file_content, content_type_for_parsing)
    metadata = {"filename": file.filename, "content_type": file.content_type}

    db_document = Document(
        user_id=current_user.id,
        filename=unique_filename,
        original_filename=file.filename,
        content=content,
        doc_metadata=metadata
    )
    db.add(db_document)
    db.commit()
    db.refresh(db_document)

    index_document(db_document.id, db_document.content, content_type_for_parsing)

    return db_document

class QueryRequest(BaseModel):
    query: str

@router.post("/query")
async def query_document(request: QueryRequest, current_user: User = Depends(get_current_user)):
    from .rag import query_documents
    docs = query_documents(request.query)
    if not docs:
        return {"answer": "No relevant documents found"}
    return {"answer": docs[0].page_content}